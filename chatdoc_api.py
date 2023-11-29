from fastapi import FastAPI, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional
import re
import uuid
from langchain.embeddings import OpenAIEmbeddings

from langchain.vectorstores import Chroma
from langchain.text_splitter import CharacterTextSplitter, RecursiveCharacterTextSplitter
from langchain.chains import ConversationalRetrievalChain
from langchain.llms import OpenAI
from langchain.document_loaders import PyPDFLoader, YoutubeLoader, PyPDFDirectoryLoader
import os
import requests
import shutil
from bs4 import BeautifulSoup
from langchain.memory import PostgresChatMessageHistory
import docx
from langchain.docstore.document import Document
from pptx import Presentation
postgres_conn = "postgresql://postgres:kclQHhW60TZkplCq@db.hcpihzivghfibozedtyu.supabase.co:5432/postgres"

llm = OpenAI(temperature=0)

def extract_text_from_pptx(pptx_file_path):
    pptx_texts = []
    presentation = Presentation(pptx_file_path)

    i=1
    for slide in presentation.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text

        pptx_texts.append(Document(page_content=text,metadata={"source":f"slide-{i}"}))
        i += 1

    return pptx_texts

def extract_text_from_docx(docx_file_path):
    docx_texts = []
    doc = docx.Document(docx_file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text
    docx_texts.append(Document(page_content=text, metadata=""))
    return docx_texts

def process_pdf(file):
    # Load the PDF file
    pdf_path = file.name
    loader = PyPDFLoader(pdf_path)
    documents = loader.load()

    # Chunk and Embeddings
    text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=0)
    documents = text_splitter.split_documents(documents)
    return documents

def process_yt(yt_link):
    loader = YoutubeLoader.from_youtube_url(yt_link)
    transcript = loader.load()
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=0)
    texts = text_splitter.split_documents(transcript)
    return texts

def crawl(url):
    output =''
    response = requests.get(url.strip())
    soup = BeautifulSoup(response.text, 'html.parser')
        # Collect the data from the webpage
    data = soup.get_text()
    output += data
        
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=0, separators=[" ", ",", "\n"])
    output = text_splitter.split_text(output)
    return output


def process_url(urls):
    yt_links = []
    web_links = []
    for i in range(len(urls)):
        txt = urls[i]
        x = re.search("(?:https?:\/\/)?(?:www\.)?youtu\.?be(?:\.com)?\/?.*(?:watch|embed)?(?:.*v=|v\/|\/)([\w\-_]+)\&?", txt)
        if x:
            processed = process_yt(urls[i])
            for proc in processed:
                yt_links.append(proc)
        else:
            processed = crawl(urls[i])
            for proc in processed:
                web_links.append(proc)
    
    return yt_links, web_links


def process_pdf(dir):
    docs = []
    loader = PyPDFDirectoryLoader(path=dir)
    documents = loader.load()
    text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=0)
    documents = text_splitter.split_documents(documents)
    
    for doc in documents:
        docs.append(doc)
    return docs


app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)



embeddings = OpenAIEmbeddings()

@app.post("/uploadfile")
async def create_upload_file(files: Optional[List[UploadFile]] = [], urls: Optional[List[str]] = []):
    user_id = str(uuid.uuid4())
    old_dir = os.curdir
    os.mkdir(old_dir + "/uploads/" + user_id)

    vectorstore = Chroma(embedding_function = embeddings, persist_directory = f"dbs/{user_id}")
    vectorstore.persist()
    
    if files == [] and urls == []:
        return HTTPException(status_code=404, detail="Give some input url or document please!")

    if urls != []:
        youtube_docs, url_texts = process_url(urls)
        if youtube_docs != []:
            doc_ids = [str(uuid.uuid5(uuid.NAMESPACE_DNS, str(doc.page_content))) for doc in youtube_docs]

            vectorstore.add_documents(youtube_docs, ids = doc_ids)
        if url_texts != []:
            text_ids = [str(uuid.uuid5(uuid.NAMESPACE_DNS, text)) for text in url_texts]
            vectorstore.add_texts(url_texts, ids = text_ids)
    for file in files:
        try:
            with open(f"uploads/{user_id}/{file.filename}", 'wb') as f:
                while contents := file.file.read(1024 * 1024):
                    f.write(contents)
        except Exception:
            return HTTPException(status_code=404, detail="There was an error uploading the file(s)")
        finally:
            file.file.close()
    if files != []:
        docs = process_pdf(f"uploads/{user_id}")
        doc_ids = [str(uuid.uuid5(uuid.NAMESPACE_DNS, str(doc.page_content))) for doc in docs]
        
    
        vectorstore.add_documents(docs, ids = doc_ids)

    shutil.rmtree(path=f"uploads/{user_id}")
    
    return {"user_id": user_id}

class Chat(BaseModel):
    message: str
    user_id: str


@app.post("/chat")
def chat(chat: Chat):
    history = PostgresChatMessageHistory(
        connection_string=postgres_conn,
        session_id=chat.user_id,
    )
    hist = [("User: " + history.messages[i].content, history.messages[i+1].content) for i in range(0, len(history.messages) - 1)]
    history_tup = [(turn[0], turn[1]) for turn in hist]
    embeddings = OpenAIEmbeddings()
    vectorstore = Chroma(embedding_function = embeddings, persist_directory = f"dbs/{chat.user_id}")
    qa = ConversationalRetrievalChain.from_llm(OpenAI(temperature=0.45), vectorstore.as_retriever(), max_tokens_limit=4097)
    response = qa({"question": chat.message, "chat_history": history_tup})
    history.add_user_message(chat.message)
    history.add_ai_message(response['answer'])    
    return {"answer": response['answer']}
    
@app.get("/")
def hello():
    return {"Hello": "API"}
